#!/usr/bin/env python3
"""Build the Nemo pitch deck as a native, editable .pptx (16:9)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy, sys

ABYSS = RGBColor(0x06, 0x22, 0x2B)
WATER = RGBColor(0x0C, 0x35, 0x41)
WATER2 = RGBColor(0x10, 0x42, 0x4F)
CORAL = RGBColor(0xFF, 0x6A, 0x3D)
CORALD = RGBColor(0xE0, 0x4E, 0x22)
STRIPE = RGBColor(0xF2, 0xF7, 0xF5)
AQUA = RGBColor(0x5B, 0xC4, 0xB4)
GLASS = RGBColor(0x8F, 0xAE, 0xAB)

W, H = Inches(13.333), Inches(7.5)
DISPLAY = "Century Gothic"
BODY = "Calibri"

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

def slide_bg(s, color=ABYSS):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color

def no_line(shape):
    shape.line.fill.background()

def box(s, x, y, w, h, fill=WATER, line=None, radius=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line:
        sp.line.color.rgb = line; sp.line.width = Pt(0.75)
    else:
        no_line(sp)
    sp.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp

def text(s, x, y, w, h, runs, size=18, color=STRIPE, font=BODY, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=None, wrap=True):
    """runs: str, or list of paragraphs; each paragraph is str or list of (txt, dict) runs."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(runs, str):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing:
            p.space_after = Pt(spacing)
        if isinstance(para, str):
            para = [(para, {})]
        for txt, st in para:
            r = p.add_run(); r.text = txt
            f = r.font
            f.size = Pt(st.get("size", size))
            f.color.rgb = st.get("color", color)
            f.name = st.get("font", font)
            f.bold = st.get("bold", bold)
            if st.get("spacing_char"):
                r.font._rPr.set('spc', str(st["spacing_char"]))
    return tb

def stripes(s, x, y, h=Inches(0.22)):
    """Clownfish three-stripe motif."""
    wds = [Inches(0.055), Inches(0.055), Inches(0.055)]
    cols = [CORAL, STRIPE, CORAL]
    hs = [h, h * 1.25, h * 0.8]
    cx = x
    for wd, c, hh in zip(wds, cols, hs):
        b = box(s, cx, y + (h * 1.25 - hh) / 2, wd, hh, fill=c, radius=0.5)
        cx += Emu(int(wd + Inches(0.045)))

def eyebrow(s, label, y=Inches(0.55)):
    stripes(s, Inches(0.75), y)
    text(s, Inches(1.12), y - Inches(0.015), Inches(9), Inches(0.35),
         [[(label.upper(), {"size": 12, "color": CORAL, "font": DISPLAY, "bold": True})]])

def headline(s, parts, y=Inches(1.05), size=32, w=Inches(11.8)):
    text(s, Inches(0.75), y, w, Inches(1.6), [parts], size=size, font=DISPLAY, bold=True)

def slideno(s, n, total=13):
    text(s, W - Inches(1.7), H - Inches(0.55), Inches(1.1), Inches(0.3),
         [[(f"{n:02d} / {total}", {"size": 10, "color": GLASS, "font": DISPLAY})]], align=PP_ALIGN.RIGHT)

def card(s, x, y, w, h, kicker, title, bodytxt, kicker_color=AQUA):
    box(s, x, y, w, h, fill=WATER, radius=0.06)
    pad = Inches(0.22)
    paras = [[(kicker.upper(), {"size": 10.5, "color": kicker_color, "font": DISPLAY, "bold": True})]]
    if title:
        paras.append([(title, {"size": 15, "color": STRIPE, "font": DISPLAY, "bold": True})])
    paras.append([(bodytxt, {"size": 11.5, "color": GLASS})])
    text(s, x + pad, y + pad, w - pad * 2, h - pad * 2, paras, spacing=6)

def statcard(s, x, y, w, h, num, label, src=None):
    box(s, x, y, w, h, fill=WATER, radius=0.06)
    pad = Inches(0.22)
    paras = [[(num, {"size": 40, "color": STRIPE, "font": DISPLAY, "bold": True})],
             [(label, {"size": 11.5, "color": GLASS})]]
    if src:
        paras.append([(src, {"size": 9, "color": GLASS})])
    text(s, x + pad, y + pad, w - pad * 2, h - pad * 2, paras, spacing=6)

def C(t): return (t, {"color": CORAL})
def A(t): return (t, {"color": AQUA})
def Wt(t): return (t, {})

# ---------------- 1 · TITLE ----------------
s = prs.slides.add_slide(BLANK); slide_bg(s)
glow = box(s, Inches(8.4), Inches(-1.6), Inches(7), Inches(7), fill=WATER, shape=MSO_SHAPE.OVAL)
glow.fill.fore_color.rgb = RGBColor(0x0A, 0x2E, 0x39)
eyebrow(s, "Confidential · July 2026", y=Inches(1.15))
text(s, Inches(0.72), Inches(1.55), Inches(8), Inches(1.9),
     [[("NEMO", {"size": 96, "color": STRIPE, "font": DISPLAY, "bold": True})]])
# big stripes next to wordmark
for i, (c, hh) in enumerate([(CORAL, 1.05), (STRIPE, 1.3), (CORAL, 0.85)]):
    box(s, Inches(4.55 + i * 0.21), Inches(2.05 + (1.3 - hh) / 2), Inches(0.12), Inches(hh), fill=c, radius=0.5)
text(s, Inches(0.75), Inches(3.45), Inches(6.6), Inches(1.6),
     [[("Neobank in a box. ", {"size": 20, "color": STRIPE, "bold": True}),
       ("Everything needed to launch and run a digital bank — deployed in a day, on any cloud or on-premise, operated by AI.",
        {"size": 20, "color": GLASS})]])
text(s, Inches(0.75), Inches(6.35), Inches(11), Inches(0.4),
     [[("CORE BANKING      LENDING      PAYMENTS      CARDS      COMPLIANCE      AI OPERATIONS",
        {"size": 11, "color": GLASS, "font": DISPLAY})]])
# simple phone mock
ph_x, ph_y, ph_w, ph_h = Inches(9.35), Inches(1.15), Inches(2.75), Inches(5.6)
box(s, ph_x, ph_y, ph_w, ph_h, fill=RGBColor(0x02, 0x16, 0x1D), line=GLASS, radius=0.18)
box(s, ph_x + Inches(0.95), ph_y + Inches(0.14), Inches(0.85), Inches(0.16), fill=RGBColor(0x0A, 0x14, 0x18), radius=0.5)
bal = box(s, ph_x + Inches(0.22), ph_y + Inches(0.55), ph_w - Inches(0.44), Inches(1.15), fill=CORAL, radius=0.14)
text(s, ph_x + Inches(0.38), ph_y + Inches(0.68), ph_w - Inches(0.75), Inches(0.95),
     [[("MAIN ACCOUNT", {"size": 7.5, "color": RGBColor(0xFF, 0xE2, 0xD6), "font": DISPLAY, "bold": True})],
      [("KES 48,250.75", {"size": 17, "color": RGBColor(0xFF, 0xFF, 0xFF), "font": DISPLAY, "bold": True})],
      [("Nemo Current ·· 4417", {"size": 8, "color": RGBColor(0xFF, 0xE2, 0xD6)})]], spacing=2)
acts = ["Send", "Pay", "Save", "Borrow"]
for i, a_ in enumerate(acts):
    bx = box(s, ph_x + Inches(0.24 + i * 0.60), ph_y + Inches(1.9), Inches(0.48), Inches(0.48), fill=WATER2, radius=0.25)
    text(s, ph_x + Inches(0.18 + i * 0.60), ph_y + Inches(2.42), Inches(0.62), Inches(0.22),
         [[(a_, {"size": 7, "color": GLASS})]], align=PP_ALIGN.CENTER)
box(s, ph_x + Inches(0.22), ph_y + Inches(2.8), ph_w - Inches(0.44), Inches(0.85), fill=WATER, radius=0.14)
text(s, ph_x + Inches(0.38), ph_y + Inches(2.92), ph_w - Inches(0.75), Inches(0.65),
     [[("Boda fund", {"size": 10, "color": STRIPE, "font": DISPLAY, "bold": True}), ("        64%", {"size": 9, "color": AQUA})],
      [("KES 32,000 of 50,000 · round-ups on", {"size": 7.5, "color": GLASS})]], spacing=3)
box(s, ph_x + Inches(0.38), ph_y + Inches(3.28), Inches(1.99), Inches(0.06), fill=RGBColor(0x2A, 0x4A, 0x52), radius=0.5)
box(s, ph_x + Inches(0.38), ph_y + Inches(3.28), Inches(1.27), Inches(0.06), fill=AQUA, radius=0.5)
txs = [("M-Pesa top-up", "+2,000.00", AQUA), ("Nairobi Water", "−850.00", GLASS), ("Round-up → Boda fund", "−12.50", GLASS)]
for i, (t1, amt, ac) in enumerate(txs):
    yy = ph_y + Inches(3.85 + i * 0.42)
    text(s, ph_x + Inches(0.26), yy, Inches(1.7), Inches(0.3), [[(t1, {"size": 8.5, "color": STRIPE})]])
    text(s, ph_x + Inches(1.7), yy, Inches(0.85), Inches(0.3), [[(amt, {"size": 8.5, "color": ac})]], align=PP_ALIGN.RIGHT)
slideno(s, 1)

# ---------------- 2 · PROBLEM ----------------
s = prs.slides.add_slide(BLANK); slide_bg(s)
eyebrow(s, "The problem")
headline(s, [Wt("Launching a digital bank still takes years and tens of millions — "),
             C("especially where it's needed most.")])
cw = Inches(3.83); gap = Inches(0.18); y0 = Inches(3.0); ch = Inches(3.2)
card(s, Inches(0.75), y0, cw, ch, "Build it yourself", "3–5 years, $10M+",
     "Core ledger, compliance, fraud, apps — before the first customer. Most attempts die in integration.")
card(s, Inches(0.75) + cw + gap, y0, cw, ch, "Buy a SaaS core", "A core is not a bank",
     "Mambu-class platforms ship a ledger, then hand you to system integrators for 6–18 months. No app, no ops, no market fit.")
card(s, Inches(0.75) + (cw + gap) * 2, y0, cw, ch, "Emerging-market reality", "Global stacks don't fit",
     "Mobile money rails, credit bureaus, local regulators and data-residency law are afterthoughts in Western platforms.")
slideno(s, 2)

# ---------------- 3 · PROVEN MODEL ----------------
s = prs.slides.add_slide(BLANK); slide_bg(s)
eyebrow(s, "The model is proven")
headline(s, [Wt("Credit-led neobanks work in emerging markets. "), A("Nobody sells the box.")])
text(s, Inches(0.75), Inches(2.55), Inches(11.8), Inches(1.1),
     [[("Fintech Farm (Leobank, Liobank, Simbank, Roarbank) proved the formula across six markets: partner with a licensed local bank, ship an app customers love, grow on AI-underwritten credit. But they operate each neobank themselves — the playbook doesn't scale beyond their own ventures.",
        {"size": 14, "color": GLASS})]])
y0 = Inches(3.9); ch = Inches(2.6)
statcard(s, Inches(0.75), y0, cw, ch, "1M+", "customers on Liobank Vietnam within two years of launch", "fintech-farm.com")
statcard(s, Inches(0.75) + cw + gap, y0, cw, ch, "6", "markets launched on one repeatable stack — Azerbaijan to Morocco", "fintech-farm.com")
statcard(s, Inches(0.75) + (cw + gap) * 2, y0, cw, ch, "0", "vendors offering that full stack as a product an operator can run themselves", "the gap Nemo fills")
slideno(s, 3)

# ---------------- 4 · PRODUCT ----------------
s = prs.slides.add_slide(BLANK); slide_bg(s)
eyebrow(s, "The product")
headline(s, [Wt("Nemo is the whole bank, "), C("in one box.")])
items = [
    ("Core", "Current, savings, wallet & term deposits, interest engines, double-entry GL (IFRS 9, maker-checker, tamper-evident audit)."),
    ("Credit", "Product factory → origination → schedules, penalties, collections, write-off. Overdrafts and BNPL included."),
    ("Payments", "Mobile-money-native transfers, P2P by phone number, bills & airtime, cards via issuer processors."),
    ("Risk & AI", "AI credit scoring, ML fraud detection with case management, AML monitoring and SAR workflows."),
    ("Compliance", "Regulatory reporting per market — central-bank returns, credit-bureau submissions, data protection."),
    ("Channels", "White-label customer app, web banking, USSD and agent network, plus the staff back-office console."),
    ("Treasury", "Float management, reconciliation against every rail, settlement with the partner bank."),
    ("Operations", "Monitoring, alerting and AI ops agents at every level — infrastructure to ledger integrity."),
]
cw4 = Inches(2.85); ch4 = Inches(1.72); y0 = Inches(2.55)
for i, (k, b) in enumerate(items):
    x = Inches(0.75) + (cw4 + Inches(0.17)) * (i % 4)
    y = y0 + (ch4 + Inches(0.17)) * (i // 4)
    box(s, x, y, cw4, ch4, fill=WATER, radius=0.07)
    text(s, x + Inches(0.18), y + Inches(0.15), cw4 - Inches(0.36), ch4 - Inches(0.3),
         [[(k.upper(), {"size": 10, "color": AQUA, "font": DISPLAY, "bold": True})],
          [(b, {"size": 9.5, "color": GLASS})]], spacing=4)
text(s, Inches(0.75), Inches(6.55), Inches(11.8), Inches(0.5),
     [[("16 event-driven microservices, one platform. The regulated back half is ", {"size": 12.5, "color": GLASS}),
       ("built and hardened today", {"size": 12.5, "color": CORAL, "bold": True}),
       (" — audited money paths, 230+ automated API tests, live-grade Kenya deployment.", {"size": 12.5, "color": GLASS})]])
slideno(s, 4)

# ---------------- 5 · THE APP ----------------
s = prs.slides.add_slide(BLANK); slide_bg(s)
eyebrow(s, "The last mile")
headline(s, [Wt("A bank customers "), C("love to hold.")])
text(s, Inches(0.75), Inches(2.4), Inches(11.8), Inches(0.7),
     [[("White-label and brand-packed per tenant: same platform, each neobank's own name, colours and products. Three moments from the concept design —",
        {"size": 13, "color": GLASS})]])
shots = [
    ("EVERYDAY MONEY", "Home", ["Balance & accounts up front", "Send · Pay bills · Save · Borrow", "Savings pot with round-up rules", "Mobile-money native top-ups"]),
    ("VIRTUAL CARDS, DAY ONE", "Cards", ["Virtual card issued at onboarding", "Freeze / unfreeze instantly", "Limits & 3-D Secure controls", "Disputes started in-app"]),
    ("CREDIT-LED GROWTH", "Loans", ["Pre-approved by Nemo AI", "Amount slider, transparent fee", "Decision in seconds, no paperwork", "Every decision logged & explainable"]),
]
y0 = Inches(3.25); ch5 = Inches(3.35)
for i, (k, ttl, bl) in enumerate(shots):
    x = Inches(0.75) + (cw + gap) * i
    box(s, x, y0, cw, ch5, fill=WATER, radius=0.07)
    bar = box(s, x, y0, cw, Inches(0.09), fill=[CORAL, STRIPE, AQUA][i], radius=0)
    paras = [[(k, {"size": 10.5, "color": [CORAL, STRIPE, AQUA][i], "font": DISPLAY, "bold": True})],
             [(ttl + " screen", {"size": 15, "color": STRIPE, "font": DISPLAY, "bold": True})]]
    for b in bl:
        paras.append([("·  " + b, {"size": 11.5, "color": GLASS})])
    text(s, x + Inches(0.25), y0 + Inches(0.3), cw - Inches(0.5), ch5 - Inches(0.55), paras, spacing=6)
slideno(s, 5)

# ---------------- 6 · HOW IT WORKS ----------------
s = prs.slides.add_slide(BLANK); slide_bg(s)
eyebrow(s, "How it works")
headline(s, [Wt("Symbiosis by design: "), A("the clownfish and the anemone.")])
text(s, Inches(0.75), Inches(2.5), Inches(11.8), Inches(1.0),
     [[("A clownfish thrives inside the anemone's protection; the anemone thrives on what the clownfish brings. The digital bank lives inside a licensed partner's regulatory shelter, and brings it customers, deposits and lending income.",
        {"size": 14, "color": GLASS})]])
bw = Inches(5.4); y0 = Inches(3.7); bh = Inches(1.85)
card(s, Inches(0.75), y0, bw, bh, "The anemone — licence holder", "Partner bank / MFI / DCP",
     "Holds the licence, the prudential capital and the settlement accounts. Keeps its core of record if it wants.", kicker_color=CORAL)
text(s, Inches(6.35), y0 + Inches(0.55), Inches(0.7), Inches(0.6),
     [[("⇄", {"size": 28, "color": CORAL, "bold": True})]], align=PP_ALIGN.CENTER)
card(s, Inches(7.15), y0, bw, bh, "The clownfish — Nemo tenant", "One tenant = one neobank",
     "Its own customers, products, ledgers, brand and reports — isolated on a shared platform. One installation runs many banks.", kicker_color=CORAL)
box(s, Inches(0.75), Inches(5.85), Inches(0.045), Inches(1.0), fill=CORAL, radius=0)
text(s, Inches(1.0), Inches(5.9), Inches(11.4), Inches(1.0),
     [[("Launching a neobank is configuration, not a project:", {"size": 12.5, "color": STRIPE, "bold": True})],
      [("1 · Provision tenant   →   2 · Pick a market pack (currency, rails, KYC, tax, regulator)   →   3 · Configure products   →   4 · Apply brand pack   →   ",
        {"size": 12.5, "color": GLASS}), ("a running, compliant bank.", {"size": 12.5, "color": CORAL, "bold": True})]], spacing=5)
slideno(s, 6)

# ---------------- 7 · DEPLOY ANYWHERE ----------------
s = prs.slides.add_slide(BLANK); slide_bg(s)
eyebrow(s, "Deploy anywhere")
headline(s, [Wt("One click to a bank — "), C("on our cloud, your cloud, or your basement.")])
y0 = Inches(3.0); ch = Inches(2.55)
card(s, Inches(0.75), y0, cw, ch, "One-command install", None,
     "The entire platform — services, databases, event bus, observability — ships as a single Kubernetes package with sane defaults. Sandbox bank in under an hour.")
card(s, Inches(0.75) + cw + gap, y0, cw, ch, "On-premise & air-gapped", None,
     "Data-residency laws in our markets make on-prem a feature, not a compromise. Offline bundles, licence activation, no phone-home required.")
card(s, Inches(0.75) + (cw + gap) * 2, y0, cw, ch, "Day-2 automated", None,
     "Zero-downtime upgrades, gated migrations, backups with tested recovery objectives, multi-node HA. Banks don't take maintenance windows.")
text(s, Inches(0.75), Inches(6.0), Inches(11.8), Inches(0.9),
     [[("Config over code — the iron rule.  ", {"size": 13, "color": STRIPE, "bold": True}),
       ("A new country is a market pack; a new savings product is product-factory data; a new brand is a theme. Nothing a regulator or product manager changes requires an engineer.",
        {"size": 13, "color": GLASS})]])
slideno(s, 7)

# ---------------- 8 · AI AT THE CORE ----------------
s = prs.slides.add_slide(BLANK); slide_bg(s)
eyebrow(s, "AI at the core")
headline(s, [Wt("The bank mostly "), A("runs itself."), Wt(" Humans set policy and review by exception.")])
bw = Inches(5.85); y0 = Inches(2.9); bh = Inches(3.0)
box(s, Inches(0.75), y0, bw, bh, fill=WATER, radius=0.06)
text(s, Inches(1.0), y0 + Inches(0.22), bw - Inches(0.5), bh - Inches(0.45),
     [[("AI DECIDES THE MONEY", {"size": 11, "color": AQUA, "font": DISPLAY, "bold": True})],
      [("·  Straight-through credit — apply → score → disburse, zero human touch in-policy, explainable declines", {"size": 12, "color": GLASS})],
      [("·  Fraud & AML — ML detection, AI triage drafting cases and SARs for officer sign-off", {"size": 12, "color": GLASS})],
      [("·  Collections — risk-based strategy, best time & channel, hardship detection", {"size": 12, "color": GLASS})]], spacing=8)
box(s, Inches(6.85), y0, bw, bh, fill=WATER, radius=0.06)
text(s, Inches(7.1), y0 + Inches(0.22), bw - Inches(0.5), bh - Inches(0.45),
     [[("AI RUNS THE SYSTEM", {"size": 11, "color": AQUA, "font": DISPLAY, "bold": True})],
      [("·  Ops agents — watch telemetry & business metrics, diagnose, execute runbooks behind approval gates", {"size": 12, "color": GLASS})],
      [("·  Customer AI — support chat with real banking actions, nudges, churn prediction", {"size": 12, "color": GLASS})],
      [("·  Governed — every model registered, drift-monitored, bias-tested, overridable, kill switch", {"size": 12, "color": GLASS})]], spacing=8)
text(s, Inches(0.75), Inches(6.2), Inches(11.8), Inches(0.6),
     [[("One decision engine logs every automated call the bank makes, with its reasons — ", {"size": 13, "color": GLASS}),
       ("automation a bank's risk committee can sign.", {"size": 13, "color": STRIPE, "bold": True})]])
slideno(s, 8)

# ---------------- 9 · WHY US ----------------
s = prs.slides.add_slide(BLANK); slide_bg(s)
eyebrow(s, "Why us")
headline(s, [Wt("We're not starting — "), C("the hard half is built.")])
text(s, Inches(0.75), Inches(2.5), Inches(11.8), Inches(0.9),
     [[("Most neobank stacks begin with the app and hit the wall at the ledger. Nemo grew the other way: it began as a bank-grade lending & deposits platform under a real regulator, and the app is the last mile.",
        {"size": 14, "color": GLASS})]])
stats = [("16", "event-driven microservices, one Go monorepo"),
         ("2×", "entry accounting — every shilling balanced, IFRS 9 provisioned, audit-trailed"),
         ("230+", "automated API tests over the money paths, plus UI test suites"),
         ("1st", "market live-grade: Kenya — CBK & credit-bureau reporting built in")]
y0 = Inches(3.6); ch9 = Inches(2.15)
for i, (n, l) in enumerate(stats):
    statcard(s, Inches(0.75) + (cw4 + Inches(0.17)) * i, y0, cw4, ch9, n, l)
text(s, Inches(0.75), Inches(6.15), Inches(11.8), Inches(0.8),
     [[("Battle scars included: hardened through functional and go-live audits of the money paths — idempotent events, tenant isolation, maker-checker controls. ",
        {"size": 12.5, "color": GLASS}), ("That discipline is the product.", {"size": 12.5, "color": STRIPE, "bold": True})]])
slideno(s, 9)

# ---------------- 10 · MARKET ----------------
s = prs.slides.add_slide(BLANK); slide_bg(s)
eyebrow(s, "The market")
headline(s, [Wt("Where banking is "), A("mobile-first and underserved.")])
y0 = Inches(2.7); ch10 = Inches(2.2)
statcard(s, Inches(0.75), y0, cw, ch10, "~1.4B", "adults worldwide still unbanked — concentrated in our target band", "World Bank Findex")
statcard(s, Inches(0.75) + cw + gap, y0, cw, ch10, "$1T+", "annual mobile-money transaction value, majority in Sub-Saharan Africa", "GSMA State of the Industry")
statcard(s, Inches(0.75) + (cw + gap) * 2, y0, cw, ch10, "6", "markets already validated for the credit-led model by the reference operator", "fintech-farm.com")
text(s, Inches(0.75), Inches(5.25), Inches(11.8), Inches(1.7),
     [[("Beachhead: East Africa.  ", {"size": 13, "color": STRIPE, "bold": True}),
       ("Kenya is live-grade today; Ethiopia — 120M people just opening to digital finance — is the first market-pack expansion. Then the proven emerging-market band: West Africa, MENA, Central & Southeast Asia.",
        {"size": 13, "color": GLASS})],
      [("Buyers:  ", {"size": 13, "color": STRIPE, "bold": True}),
       ("tier-2/3 banks going digital · MFIs and digital credit providers upgrading to bank-grade · licensed fintechs & telcos · neobank venture operators who want the box without building it.",
        {"size": 13, "color": GLASS})]], spacing=10)
slideno(s, 10)

# ---------------- 11 · BUSINESS MODEL ----------------
s = prs.slides.add_slide(BLANK); slide_bg(s)
eyebrow(s, "Business model & landscape")
headline(s, [Wt("Licence the box. "), C("Share the upside.")])
text(s, Inches(0.75), Inches(2.35), Inches(11.8), Inches(0.75),
     [[("Per-tenant platform licence (cloud subscription or on-prem term), metered by active customers and modules — with an optional revenue-share on lending income that aligns us with the operator.",
        {"size": 13, "color": GLASS})]])
rows = [
    ["", "Nemo", "Fintech Farm", "SaaS cores", "Legacy cores"],
    ["Full stack incl. consumer app", "✔", "✔ self-operated only", "✖ core only + SI", "partial"],
    ["On-premise / data residency", "✔", "✖", "limited", "✔ heavy"],
    ["Time to launch", "days–weeks", "months, theirs to run", "6–18 months", "18 months+"],
    ["AI-native operations", "✔ core design", "credit only", "✖", "✖"],
    ["Emerging-market packs", "✔ built in", "✔ bespoke", "per-project", "✖"],
]
tw, th = Inches(11.8), Inches(3.1)
gtbl = s.shapes.add_table(len(rows), 5, Inches(0.75), Inches(3.3), tw, th).table
gtbl.columns[0].width = Inches(3.4)
for c in range(1, 5):
    gtbl.columns[c].width = Inches(2.1)
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = gtbl.cell(ri, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = WATER if ri == 0 else (WATER2 if ci == 1 else ABYSS)
        tfc = cell.text_frame; tfc.word_wrap = True
        p = tfc.paragraphs[0]; r = p.add_run(); r.text = val
        f = r.font; f.size = Pt(11); f.name = BODY
        if ri == 0:
            f.color.rgb = AQUA; f.bold = True; f.name = DISPLAY; f.size = Pt(10.5)
        elif ci == 0:
            f.color.rgb = STRIPE; f.bold = True
        elif ci == 1:
            f.color.rgb = CORAL; f.bold = True
        else:
            f.color.rgb = GLASS
text(s, Inches(0.75), Inches(6.7), Inches(11.8), Inches(0.6),
     [[("The wedge: ", {"size": 12.5, "color": STRIPE, "bold": True}),
       ("the only neobank-complete platform an emerging-market operator can run themselves, on their own infrastructure, with AI doing the operational heavy lifting.",
        {"size": 12.5, "color": GLASS})]])
slideno(s, 11)

# ---------------- 12 · ROADMAP ----------------
s = prs.slides.add_slide(BLANK); slide_bg(s)
eyebrow(s, "Roadmap")
headline(s, [Wt("Package it. Show it. "), A("Let it run itself.")])
phases = [
    ("PHASE 1 · THIS QUARTER", "Package", CORAL,
     ["One-command Kubernetes install", "Tenant provisioning: create-neobank end-to-end", "Market packs — Kenya extracted as pack #1",
      "Observability: tracing, business alerting, SLOs", "Self-service eKYC onboarding APIs", "Explainable decision engine v1"]),
    ("PHASE 2 · QUARTERS 2–3", "The visible neobank", STRIPE,
     ["White-label customer app: accounts, P2P, bills, loans", "Virtual cards via issuer processor", "Savings pots, scheduled payments, deposit lifecycle",
      "Reconciliation engine across all rails", "Straight-through credit + model governance", "On-prem installer, zero-downtime upgrades, HA/DR"]),
    ("PHASE 3 · QUARTER 4+", "AI-operated & ecosystem", AQUA,
     ["Ops agents: diagnose & remediate behind gates", "AI collections, AML copilot, customer AI", "USSD + agent network channel",
      "Market pack #2: Ethiopia — config-over-code proven", "Public API platform & developer portal", "Tenant metering & platform billing"]),
]
y0 = Inches(2.6); ch12 = Inches(3.75)
for i, (when, ttl, accent, bl) in enumerate(phases):
    x = Inches(0.75) + (cw + gap) * i
    box(s, x, y0, cw, ch12, fill=WATER, radius=0.05)
    box(s, x, y0, cw, Inches(0.08), fill=accent, radius=0)
    paras = [[(when, {"size": 9.5, "color": GLASS, "font": DISPLAY, "bold": True})],
             [(ttl, {"size": 15, "color": STRIPE, "font": DISPLAY, "bold": True})]]
    for b in bl:
        paras.append([("·  " + b, {"size": 10.5, "color": GLASS})])
    text(s, x + Inches(0.22), y0 + Inches(0.25), cw - Inches(0.44), ch12 - Inches(0.5), paras, spacing=5)
text(s, Inches(0.75), Inches(6.6), Inches(11.8), Inches(0.5),
     [[("Standing tracks: money-path audits · security posture · PCI-DSS → ISO 27001 → SOC 2 certification programme.",
        {"size": 11.5, "color": GLASS})]])
slideno(s, 12)

# ---------------- 13 · CLOSE ----------------
s = prs.slides.add_slide(BLANK); slide_bg(s)
eyebrow(s, "Nemo", y=Inches(1.5))
text(s, Inches(0.75), Inches(2.0), Inches(11.2), Inches(2.4),
     [[("Every market that needs a bank should be ", {"size": 40, "color": STRIPE, "font": DISPLAY, "bold": True}),
       ("one click away", {"size": 40, "color": CORAL, "font": DISPLAY, "bold": True}),
       (" from one.", {"size": 40, "color": STRIPE, "font": DISPLAY, "bold": True})]])
text(s, Inches(0.75), Inches(4.6), Inches(10.5), Inches(1.0),
     [[("The regulated core is built. The model is proven. What remains is packaging and the last mile to the customer's hand — and that is exactly what the next three quarters deliver.",
        {"size": 16, "color": GLASS})]])
text(s, Inches(0.75), Inches(6.3), Inches(11), Inches(0.4),
     [[("NEMO · NEOBANK IN A BOX          NAIROBI · JULY 2026", {"size": 11, "color": GLASS, "font": DISPLAY})]])
slideno(s, 13)

out = sys.argv[1] if len(sys.argv) > 1 else "nemo-pitch-deck.pptx"
prs.save(out)
print("saved", out)
